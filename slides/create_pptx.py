from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
CORAL = RGBColor(0xD4, 0x65, 0x4A)
CORAL_LIGHT = RGBColor(0xE8, 0x85, 0x6D)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
NAVY2 = RGBColor(0x22, 0x22, 0x3A)
WARM = RGBColor(0xFA, 0xF8, 0xF5)
WARM2 = RGBColor(0xF2, 0xEF, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x18, 0x18, 0x1B)
TEXT_BODY = RGBColor(0x3A, 0x3A, 0x42)
TEXT_MUTED = RGBColor(0x7A, 0x7A, 0x88)
ACCENT = RGBColor(0xC9, 0xA9, 0x6E)
BORDER = RGBColor(0xE4, 0xE0, 0xDA)
CORAL_BG = RGBColor(0xFB, 0xF0, 0xEC)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

SLIDE_W = Inches(16)
SLIDE_H = Inches(9)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    else:
        shape.adjustments[0] = 0.02
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_BODY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Noto Sans JP'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_multiline(slide, left, top, width, height, lines, font_size=18, color=TEXT_BODY, bold=False, line_spacing=1.8, alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Noto Sans JP'
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txbox


def add_bullet_list(slide, left, top, width, height, items, font_size=20, color=TEXT_BODY, bullet_color=CORAL):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Noto Sans JP'
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.level = 0
        # bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u25A0')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgb = etree.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', '%02X%02X%02X' % (bullet_color[0], bullet_color[1], bullet_color[2]))
        buSz = etree.SubElement(pPr, qn('a:buSzPct'))
        buSz.set('val', '60000')
    return txbox


def add_label(slide, left, top, text, bg_color=None, text_color=CORAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.5), Inches(0.4))
    shape.fill.background()
    shape.line.color.rgb = text_color
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.3
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = 'Inter'
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_highlight_box(slide, left, top, width, height, text, font_size=22):
    shape = add_shape_rect(slide, left, top, width, height, CORAL_BG, CORAL)
    shape.adjustments[0] = 0.04
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = TEXT_DARK
    p.font.bold = True
    p.font.name = 'Noto Sans JP'
    p.alignment = PP_ALIGN.LEFT
    return shape


def add_card(slide, left, top, width, height, title, desc, title_color=TEXT_DARK):
    shape = add_shape_rect(slide, left, top, width, height, WHITE, BORDER)
    shape.adjustments[0] = 0.04
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.35)
    tf.margin_bottom = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = 'Noto Sans JP'
    p.space_after = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Noto Sans JP'
    return shape


def add_stat_card(slide, left, top, width, height, number, source, desc):
    shape = add_shape_rect(slide, left, top, width, height, WHITE, BORDER)
    shape.adjustments[0] = 0.04
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.35)
    tf.margin_bottom = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(44)
    p.font.color.rgb = CORAL
    p.font.bold = True
    p.font.name = 'Inter'
    p.space_after = Pt(6)
    p2 = tf.add_paragraph()
    p2.text = source
    p2.font.size = Pt(11)
    p2.font.color.rgb = CORAL
    p2.font.bold = True
    p2.font.name = 'Inter'
    p2.space_after = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = 'Noto Sans JP'
    return shape


def add_flow_step(slide, left, top, text):
    shape = add_shape_rect(slide, left, top, Inches(1.8), Inches(1.2), WHITE, BORDER)
    shape.adjustments[0] = 0.06
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_DARK
    p.font.bold = True
    p.font.name = 'Noto Sans JP'
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_flow_arrow(slide, left, top):
    add_text(slide, left, top, Inches(0.5), Inches(1.2), '→', font_size=24, color=CORAL, bold=True, alignment=PP_ALIGN.CENTER)


def section_slide(title, subtitle='', label_text=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, NAVY)
    # Decorative accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(0.5), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = CORAL
    line.line.fill.background()
    if label_text:
        add_label(slide, Inches(6.75), Inches(2.8), label_text, text_color=CORAL_LIGHT)
    add_text(slide, Inches(1.5), Inches(3.4), Inches(13), Inches(2.5), title, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1.5), Inches(5.8), Inches(13), Inches(1), subtitle, font_size=22, color=RGBColor(0x99, 0x99, 0xAA), bold=True, alignment=PP_ALIGN.CENTER)
    return slide


def content_slide(bg_color=WARM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg_color)
    return slide


# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
s = section_slide('Claude Code 導入勉強会', 'コピペ地獄から解放される AI実務活用の新常識', '2026.4.10  ONLINE SEMINAR')

# =====================================================================
# WHY THIS SEMINAR
# =====================================================================
section_slide('なぜこの勉強会を\nやるのか', '', 'PROLOGUE')

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), 'BACKGROUND')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1.5),
         '行政書士として\nコピペ地獄に苦しんでいた', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(3.5), [
    '役所の書式にデータを手入力する毎日',
    'Claude Codeを業務に組み込んだことで、実務が根本的に変わった',
    'しかし情報はエンジニア向けか汎用的なものばかり',
    '士業・小規模事業者の実務に落とし込んだ情報がどこにもない',
], font_size=20)
add_highlight_box(s, Inches(1.2), Inches(6.5), Inches(13), Inches(1.0),
                  'だから自分でやることにしました', font_size=24)

# =====================================================================
# INTRODUCTION
# =====================================================================
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), 'INTRODUCTION')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1), '自己紹介', font_size=36, color=TEXT_DARK, bold=True)
add_text(s, Inches(1.2), Inches(2.5), Inches(13), Inches(1.5),
         '行政書士 阿久津和宏\n補助金申請の実務にClaude Codeを導入している当事者', font_size=22, color=TEXT_BODY)
add_highlight_box(s, Inches(1.2), Inches(4.8), Inches(13), Inches(2.0),
                  '今日のルール\nプログラミングの話はしません。コードは1行も書きません。', font_size=24)

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), 'こんな経験ありませんか？')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1.5),
         '今日の参加者が抱えている痛み', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(4.5), [
    'AIやツールを入れたのに帰宅時間が変わらない',
    '「自動化できる」と聞いたのに自動化する作業に追われている',
    'SNSを見ていると乗り遅れを感じて焦る',
    '役所のポンコツ書式のせいでAIが使えない',
], font_size=22)

# =====================================================================
# PART 1
# =====================================================================
section_slide('なぜAIを入れても\n楽にならないのか', '', 'PART 1')

# Loop
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), 'あるある')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1), 'AIツール導入の無限ループ', font_size=36, color=TEXT_DARK, bold=True)
x_start = Inches(1.5)
steps = ['ChatGPT\n契約', 'ツール\n試す', 'YouTube\n勉強', '楽に\nならない', 'また\n探す']
for i, step in enumerate(steps):
    add_flow_step(s, x_start + Inches(i * 2.7), Inches(3.5), step)
    if i < len(steps) - 1:
        add_flow_arrow(s, x_start + Inches(i * 2.7 + 1.8), Inches(3.5))
add_text(s, Inches(1.2), Inches(6.5), Inches(13.5), Inches(1),
         'ツールを使うための新たな手間が増えている', font_size=22, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# Data
s = content_slide(WHITE)
add_label(s, Inches(1.2), Inches(0.6), 'DATA')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '個人の能力ではなく構造の問題', font_size=36, color=TEXT_DARK, bold=True)
add_stat_card(s, Inches(1.2), Inches(3.2), Inches(6.5), Inches(4.0),
              '77%', 'Upwork調査', 'AI導入現場の77%が\n逆に仕事量が増加している')
add_stat_card(s, Inches(8.3), Inches(3.2), Inches(6.5), Inches(4.0),
              '80%', 'Pendo調査', 'クラウドツール機能の80%は\n実際には使われていない')

# Data carrier
s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '正体')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'あなたは「データの運び屋」になっている', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(2.5), [
    'PDFを開く → 別システムに手打ち',
    'ChatGPTの出力 → コピー → 役所フォーマットに貼り付け',
    'ログイン → パスワード → 認証コード → やっとスタート',
], font_size=22)
add_highlight_box(s, Inches(1.2), Inches(6.0), Inches(13), Inches(1.5),
                  'どれも「専門知識」を使っていない作業\nツールとツールの間を人間が手作業で埋めている', font_size=22)

# Time loss
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '時間のロス')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '「たかがコピペ」が奪っているもの', font_size=36, color=TEXT_DARK, bold=True)
add_stat_card(s, Inches(1.2), Inches(3.2), Inches(6.5), Inches(4.0),
              '23分15秒', 'カリフォルニア大学', '作業を切り替えてから\n集中に戻るまでにかかる時間')
add_stat_card(s, Inches(8.3), Inches(3.2), Inches(6.5), Inches(4.0),
              '約5週間/年', 'HBR', 'アプリ切り替えだけで\n年間約5週間分のロス')

# Brain damage
s = content_slide(WHITE)
add_label(s, Inches(1.2), Inches(0.6), '脳と家庭へのダメージ')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'コピペは脳をすり減らしている', font_size=36, color=TEXT_DARK, bold=True)
add_card(s, Inches(1.2), Inches(3.2), Inches(6.5), Inches(4.0),
         '脳の疲労', '繰り返し認知タスクで疲労物質が蓄積（パリ脳研究所）。\n現状維持バイアスが強化され、新規事業に踏み出せなくなる')
add_card(s, Inches(8.3), Inches(3.2), Inches(6.5), Inches(4.0),
         '家庭への波及', '未完了タスクが不眠を誘発（ベイラー大学）。\n前頭葉疲労で感情制御が低下し、家庭へ持ち込まれる')

# =====================================================================
# PART 2
# =====================================================================
section_slide('AI情報を追いかけても\n自分の作業が変わらない理由', '', 'PART 2')

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '問題 1')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'ツール紹介で止まっている', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(4.5), [
    '世の中のAI情報の大半は「このツールが便利」で終わっている',
    'それを自分の業務のどこに・どう組み込むかは誰も教えてくれない',
    'ツールを知ることと、業務が楽になることは別の話',
], font_size=22)

s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '問題 2')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '汎用的すぎて実務に合わない', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(2.5), [
    '「議事録要約」「メール文案」→ 業務の末端で本丸じゃない',
    '本当に時間を取られるのは役所の独自書式、申請手続き、顧客別の書類',
    '発信者は「AIの専門家」であって「あなたの業務の専門家」ではない',
], font_size=22)
add_highlight_box(s, Inches(1.2), Inches(6.0), Inches(13), Inches(1.5),
                  'あなたの業務を知らない人の情報は\nどれだけ正しくても現場には刺さらない', font_size=22)

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '問題 3')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '「全部AIにやらせる」という誤解', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(2.0), [
    'AIに丸投げ → 出力を修正 → 修正作業が増える',
    '「コピペの運び屋」が「修正の運び屋」に変わっただけ',
], font_size=22)
add_highlight_box(s, Inches(1.2), Inches(5.5), Inches(13), Inches(1.5),
                  'AIに任せる部分と自分がやる部分の線引きは\n業務の中身を知っている人間にしかできない', font_size=22)

# =====================================================================
# PART 3
# =====================================================================
section_slide('AI導入の\nマインドセット', '', 'PART 3')

# Over/Under estimation
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '過大評価 or 過小評価')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'AIの能力を正しく知る', font_size=36, color=TEXT_DARK, bold=True)
add_card(s, Inches(1.2), Inches(3.2), Inches(6.5), Inches(3.0),
         '過大評価すると...', '「丸投げしたのに使えない」\nすぐ諦める', title_color=TEXT_MUTED)
add_card(s, Inches(8.3), Inches(3.2), Inches(6.5), Inches(3.0),
         '過小評価すると...', '「どうせ自分の業務には合わない」\n触らない', title_color=TEXT_MUTED)
add_text(s, Inches(1.2), Inches(7.0), Inches(13.5), Inches(1),
         'この境界線を知らないまま使うから「期待はずれ」になる', font_size=22, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# Can / Cannot
s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '境界線')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'AIにできること・できないこと', font_size=36, color=TEXT_DARK, bold=True)
# できること
shape_ok = add_shape_rect(s, Inches(1.2), Inches(3.2), Inches(6.5), Inches(4.5), CORAL_BG, CORAL)
shape_ok.adjustments[0] = 0.04
add_text(s, Inches(1.6), Inches(3.4), Inches(5.5), Inches(0.4), 'AIにできること', font_size=13, color=CORAL, bold=True)
add_bullet_list(s, Inches(1.6), Inches(4.0), Inches(5.8), Inches(3.5), [
    '大量のデータ処理', '定型文の生成', '情報の収集・整理', 'ファイル操作の自動実行',
], font_size=18)
# できないこと
shape_ng = add_shape_rect(s, Inches(8.3), Inches(3.2), Inches(6.5), Inches(4.5), RGBColor(0xF5, 0xF3, 0xEF), BORDER)
shape_ng.adjustments[0] = 0.04
add_text(s, Inches(8.7), Inches(3.4), Inches(5.5), Inches(0.4), 'AIにできないこと', font_size=13, color=TEXT_MUTED, bold=True)
add_bullet_list(s, Inches(8.7), Inches(4.0), Inches(5.8), Inches(3.5), [
    'クライアントの本音を汲む', '事業の本質的な強みを見抜く', '審査基準との最終判断',
], font_size=18, bullet_color=TEXT_MUTED)

# 切り分け
s = content_slide(WHITE)
add_label(s, Inches(1.2), Inches(0.6), '補助金 事業計画書の場合')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '正しい切り分けの例', font_size=36, color=TEXT_DARK, bold=True)
shape_ai = add_shape_rect(s, Inches(1.2), Inches(3.0), Inches(6.5), Inches(3.2), CORAL_BG, CORAL)
shape_ai.adjustments[0] = 0.04
add_text(s, Inches(1.6), Inches(3.2), Inches(5.5), Inches(0.4), 'AIが担当', font_size=13, color=CORAL, bold=True)
add_bullet_list(s, Inches(1.6), Inches(3.8), Inches(5.8), Inches(2.2), [
    '市場データの収集', '統計の整理', '文章の構成・初稿',
], font_size=18)
shape_me = add_shape_rect(s, Inches(8.3), Inches(3.0), Inches(6.5), Inches(3.2), RGBColor(0xF5, 0xF3, 0xEF), BORDER)
shape_me.adjustments[0] = 0.04
add_text(s, Inches(8.7), Inches(3.2), Inches(5.5), Inches(0.4), '自分が担当', font_size=13, color=TEXT_MUTED, bold=True)
add_bullet_list(s, Inches(8.7), Inches(3.8), Inches(5.8), Inches(2.2), [
    'ヒアリング', '事業の強みの言語化', '審査基準との最終照合',
], font_size=18, bullet_color=TEXT_MUTED)
add_highlight_box(s, Inches(1.2), Inches(6.8), Inches(13.5), Inches(1.2),
                  'この切り分けができるとコピペが消える\nAIが直接ファイルを読み、直接書き出す', font_size=22)

# Who does it
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '結論')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1.5),
         'じゃあその切り分けは\n誰がやるのか', font_size=36, color=TEXT_DARK, bold=True)
add_highlight_box(s, Inches(1.2), Inches(3.8), Inches(13), Inches(1.2),
                  '業務の中身を知っている人間にしかできない', font_size=28)
add_text(s, Inches(1.2), Inches(5.8), Inches(13), Inches(1.5),
         'AIの専門家ではなく\n業務の専門家がAIの組み込み方を決める', font_size=22, color=TEXT_BODY)

# =====================================================================
# PART 4
# =====================================================================
section_slide('AI活用にかかせない\n3つの要素', '', 'PART 4')

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '3 ELEMENTS')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '揃えるべき3つの要素', font_size=36, color=TEXT_DARK, bold=True)
add_card(s, Inches(1.2), Inches(3.2), Inches(4.2), Inches(3.5),
         '役割（誰が）', 'どのAIに、どの作業を担当させるか', title_color=CORAL)
add_card(s, Inches(5.9), Inches(3.2), Inches(4.2), Inches(3.5),
         '指示（何を）', '業務固有の手順・書式・基準をどう伝えるか', title_color=CORAL)
add_card(s, Inches(10.6), Inches(3.2), Inches(4.2), Inches(3.5),
         '参考（どのように）', 'どんな素材・データ・前例を渡すか', title_color=CORAL)

s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '3つが揃うと')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '何が変わるのか', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(3.0), [
    '役割が決まっている → AIが勝手に動く',
    '指示が明確 → 出力の手直しが激減する',
    '参考が揃っている → 自分の業務に合った出力が出る',
], font_size=22)
add_highlight_box(s, Inches(1.2), Inches(6.0), Inches(13), Inches(1.2),
                  '1つでも欠けると「便利だけど楽にならない」が続く', font_size=22)

s = content_slide(WHITE)
add_label(s, Inches(1.2), Inches(0.6), '具体例')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '3要素を実務に当てはめる', font_size=36, color=TEXT_DARK, bold=True)
add_card(s, Inches(1.2), Inches(3.2), Inches(4.2), Inches(4.0),
         '役割', '「あなたは補助金申請の専門家で、事業計画書の初稿を作成する担当です」', title_color=CORAL)
add_card(s, Inches(5.9), Inches(3.2), Inches(4.2), Inches(4.0),
         '指示', '「この書式に沿って、審査基準5項目を全て網羅した計画書を書いてください」', title_color=CORAL)
add_card(s, Inches(10.6), Inches(3.2), Inches(4.2), Inches(4.0),
         '参考', '「ヒアリングシート、過去の採択事例、市場データを素材として使ってください」', title_color=CORAL)

# =====================================================================
# PART 5
# =====================================================================
section_slide('Claude Code\n導入マニュアル解説', '', 'PART 5')

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), 'ChatGPTとの違い')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'Claude Codeとは何か', font_size=36, color=TEXT_DARK, bold=True)
shape_chat = add_shape_rect(s, Inches(1.2), Inches(3.0), Inches(6.5), Inches(4.5), RGBColor(0xF5, 0xF3, 0xEF), BORDER)
shape_chat.adjustments[0] = 0.04
add_text(s, Inches(1.6), Inches(3.2), Inches(5.5), Inches(0.4), 'チャットAI', font_size=13, color=TEXT_MUTED, bold=True)
add_bullet_list(s, Inches(1.6), Inches(3.8), Inches(5.8), Inches(3.5), [
    '質問すると答えを返す', 'コピペは自分', 'ツール間の橋渡しは人間',
], font_size=18, bullet_color=TEXT_MUTED)
shape_cc = add_shape_rect(s, Inches(8.3), Inches(3.0), Inches(6.5), Inches(4.5), CORAL_BG, CORAL)
shape_cc.adjustments[0] = 0.04
add_text(s, Inches(8.7), Inches(3.2), Inches(5.5), Inches(0.4), 'Claude Code', font_size=13, color=CORAL, bold=True)
add_bullet_list(s, Inches(8.7), Inches(3.8), Inches(5.8), Inches(3.5), [
    '指示を出せばAIが直接ファイルを読み書き', 'コピペが発生しない', '「運び屋」が構造的に消える',
], font_size=18)

# CLAUDE.md
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '心臓部')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'CLAUDE.md = パート4の3要素そのもの', font_size=36, color=TEXT_DARK, bold=True)
add_text(s, Inches(1.2), Inches(2.5), Inches(13), Inches(0.8),
         'Claude Codeの心臓部にCLAUDE.mdという設定ファイルがあります', font_size=20, color=TEXT_BODY)
add_card(s, Inches(1.2), Inches(3.8), Inches(4.2), Inches(3.0),
         '役割', 'あなたは○○の専門家で、○○の業務を担当する', title_color=CORAL)
add_card(s, Inches(5.9), Inches(3.8), Inches(4.2), Inches(3.0),
         '指示', 'この手順で、この書式で、この基準に沿って作業しろ', title_color=CORAL)
add_card(s, Inches(10.6), Inches(3.8), Inches(4.2), Inches(3.0),
         '参考', 'この素材・データ・前例を使え', title_color=CORAL)
add_text(s, Inches(1.2), Inches(7.3), Inches(13.5), Inches(0.6),
         'これがあるからAIが毎回文脈を理解した状態で動く', font_size=20, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# Install flow
s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), 'LIVE DEMO')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '画面で導入手順を見せます', font_size=36, color=TEXT_DARK, bold=True)
flow_steps = ['インストール', '認証', 'CLAUDE.md\n作成・配置', 'タスク実行']
for i, step in enumerate(flow_steps):
    add_flow_step(s, Inches(1.5 + i * 3.5), Inches(3.5), step)
    if i < len(flow_steps) - 1:
        add_flow_arrow(s, Inches(1.5 + i * 3.5 + 1.8), Inches(3.5))
add_highlight_box(s, Inches(1.2), Inches(6.0), Inches(13), Inches(1.5),
                  '今見せた手順は参加特典のマニュアルにすべて入っています\n帰ったらそのまま再現できます', font_size=22)

# =====================================================================
# PART 6 BONUS
# =====================================================================
section_slide('事業計画書を\n一晩で作成する方法', '', 'PART 6 — BONUS')

s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '前振り')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1.5),
         '3要素 + CLAUDE.mdが\n業務に組み込まれるとどうなるか', font_size=36, color=TEXT_DARK, bold=True)
add_text(s, Inches(1.2), Inches(3.5), Inches(13), Inches(1.5),
         'パート4で話した3要素、パート5で見せたCLAUDE.md。\nこれが実際の業務でどう動くか、ここからデモで見せます。', font_size=22, color=TEXT_BODY)

# Demo flow
s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), 'LIVE DEMO')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'ライブデモ', font_size=36, color=TEXT_DARK, bold=True)
add_text(s, Inches(1.2), Inches(2.3), Inches(13), Inches(0.6),
         '一言指示を入れるだけで、スキルが自動連携して動きます', font_size=20, color=TEXT_BODY)
demo_steps = ['ヒアリング\n整理', '市場調査', 'SWOT\n分析', '計画書\n作成', '審査\nシミュレーション']
for i, step in enumerate(demo_steps):
    add_flow_step(s, Inches(0.8 + i * 2.9), Inches(3.8), step)
    if i < len(demo_steps) - 1:
        add_flow_arrow(s, Inches(0.8 + i * 2.9 + 1.8), Inches(3.8))
add_highlight_box(s, Inches(1.2), Inches(6.5), Inches(13), Inches(1.0),
                  '各ステップが自動で次に進む = 人間のコピペが一切ない', font_size=22)

# Why overnight
s = content_slide(WHITE)
add_label(s, Inches(1.2), Inches(0.6), 'なぜ「一晩」なのか')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'なぜ「一晩」なのか', font_size=36, color=TEXT_DARK, bold=True)
shape_before = add_shape_rect(s, Inches(1.2), Inches(3.0), Inches(6.5), Inches(4.5), RGBColor(0xF5, 0xF3, 0xEF), BORDER)
shape_before.adjustments[0] = 0.04
add_text(s, Inches(1.6), Inches(3.2), Inches(5.5), Inches(0.4), '従来', font_size=13, color=TEXT_MUTED, bold=True)
add_bullet_list(s, Inches(1.6), Inches(3.8), Inches(5.8), Inches(3.5), [
    'ヒアリング', '手動で市場調査', 'Excelで分析', 'Wordで作文', '何日もかかる',
], font_size=18, bullet_color=TEXT_MUTED)
shape_after = add_shape_rect(s, Inches(8.3), Inches(3.0), Inches(6.5), Inches(4.5), CORAL_BG, CORAL)
shape_after.adjustments[0] = 0.04
add_text(s, Inches(8.7), Inches(3.2), Inches(5.5), Inches(0.4), 'Claude Code', font_size=13, color=CORAL, bold=True)
add_bullet_list(s, Inches(8.7), Inches(3.8), Inches(5.8), Inches(3.5), [
    '夕食前に情報を入力して走らせる', '食べている間にノンストップで進む', '戻ったら初稿が完成している',
], font_size=18)

# Same frame
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '気づき')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '同じフレームで動いている', font_size=36, color=TEXT_DARK, bold=True)
shape_plan = add_shape_rect(s, Inches(1.2), Inches(3.0), Inches(6.5), Inches(3.0), CORAL_BG, CORAL)
shape_plan.adjustments[0] = 0.04
add_text(s, Inches(1.6), Inches(3.2), Inches(5.5), Inches(0.4), '事業計画書の構造', font_size=13, color=CORAL, bold=True)
add_bullet_list(s, Inches(1.6), Inches(3.8), Inches(5.8), Inches(2.0), [
    '誰が（事業者）', '何を（補助事業）', 'どのように（実施方法・根拠データ）',
], font_size=18)
shape_ai2 = add_shape_rect(s, Inches(8.3), Inches(3.0), Inches(6.5), Inches(3.0), CORAL_BG, CORAL)
shape_ai2.adjustments[0] = 0.04
add_text(s, Inches(8.7), Inches(3.2), Inches(5.5), Inches(0.4), 'AIへの指示の構造', font_size=13, color=CORAL, bold=True)
add_bullet_list(s, Inches(8.7), Inches(3.8), Inches(5.8), Inches(2.0), [
    '役割（誰が）', '指示（何を）', '参考（どのように）',
], font_size=18)
add_highlight_box(s, Inches(1.2), Inches(6.8), Inches(13.5), Inches(1.0),
                  'AIへの指示も事業計画書も、同じフレームで動いている', font_size=22)

# =====================================================================
# PART 7 BONUS
# =====================================================================
section_slide('補助金自動収集\nシステムのデモ', '', 'PART 7 — BONUS')

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '現状')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '今やっていること', font_size=36, color=TEXT_DARK, bold=True)
manual_steps = ['各省庁\nサイト巡回', '情報を\n読む', '関係あるか\n判断', '別の場所に\nメモ']
for i, step in enumerate(manual_steps):
    add_flow_step(s, Inches(1.5 + i * 3.5), Inches(3.5), step)
    if i < len(manual_steps) - 1:
        add_flow_arrow(s, Inches(1.5 + i * 3.5 + 1.8), Inches(3.5))
add_text(s, Inches(1.2), Inches(6.5), Inches(13.5), Inches(1),
         'まさにコピペ作業。しかも毎日やらないと見落とす。', font_size=22, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), 'LIVE DEMO')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '自動収集システム', font_size=36, color=TEXT_DARK, bold=True)
auto_steps = ['毎日\n自動収集', 'フィルタリング', '通知']
for i, step in enumerate(auto_steps):
    add_flow_step(s, Inches(2.5 + i * 4.0), Inches(3.5), step)
    if i < len(auto_steps) - 1:
        add_flow_arrow(s, Inches(2.5 + i * 4.0 + 1.8), Inches(3.5))
add_text(s, Inches(1.2), Inches(6.5), Inches(13.5), Inches(1),
         '人間がやっていた巡回・転記が丸ごと消えています', font_size=22, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

s = content_slide(WHITE)
add_label(s, Inches(1.2), Inches(0.6), 'パート1からの回収')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'これがコピペが消えた状態', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(4.0), [
    'パート1からずっと話してきた「運び屋」がここでは存在しない',
    'サイトを開く必要もない、メモする必要もない、見落としもない',
    'これがパート4の3要素が揃った状態の、一番身近な例',
], font_size=22)

# =====================================================================
# OFFER
# =====================================================================
section_slide('補助金攻略コースの\nご案内', '', 'OFFER')

s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '転換')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'ここまで見せてきたこと', font_size=36, color=TEXT_DARK, bold=True)
add_bullet_list(s, Inches(1.2), Inches(3.0), Inches(13), Inches(2.5), [
    'コピペが消える仕組み',
    '3要素の考え方',
    'Claude Codeの導入',
    '事業計画書と補助金収集の実際のデモ',
], font_size=22)
add_highlight_box(s, Inches(1.2), Inches(6.0), Inches(13), Inches(1.5),
                  'これを補助金申請に全部組み込んで\n一発で動くようにしたコースがあります', font_size=22)

# Course contents
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), 'コースの中身')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'コース内容', font_size=36, color=TEXT_DARK, bold=True)
add_card(s, Inches(1.2), Inches(3.2), Inches(4.2), Inches(4.0),
         '動画講座 40本超', '補助金の基礎から審査対策、電子申請まで網羅', title_color=CORAL)
add_card(s, Inches(5.9), Inches(3.2), Inches(4.2), Inches(4.0),
         'AIスキル 5本', 'ヒアリング整理、市場調査、SWOT分析、計画書作成、審査シミュレーション', title_color=CORAL)
add_card(s, Inches(10.6), Inches(3.2), Inches(4.2), Inches(4.0),
         'テンプレート集', '採択事例データベース付き', title_color=CORAL)

# Safety
s = content_slide(WARM)
add_label(s, Inches(1.2), Inches(0.6), '安心サポート')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '3つの安心', font_size=36, color=TEXT_DARK, bold=True)
add_card(s, Inches(1.2), Inches(3.2), Inches(4.2), Inches(3.5),
         '30日間全額返金保証', '合わなければ全額お返しします', title_color=CORAL)
add_card(s, Inches(5.9), Inches(3.2), Inches(4.2), Inches(3.5),
         'メールサポート', '困ったときはいつでも相談できます', title_color=CORAL)
add_card(s, Inches(10.6), Inches(3.2), Inches(4.2), Inches(3.5),
         '永久アップデート', '制度変更にも追加費用なしで対応', title_color=CORAL)

# Price
s = content_slide(WHITE)
add_label(s, Inches(1.2), Inches(0.6), '価格比較')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '価格と比較', font_size=36, color=TEXT_DARK, bold=True)
# Old price
shape_old = add_shape_rect(s, Inches(1.2), Inches(3.0), Inches(6.5), Inches(4.5), RGBColor(0xF5, 0xF3, 0xEF), BORDER)
shape_old.adjustments[0] = 0.04
add_text(s, Inches(1.2), Inches(3.4), Inches(6.5), Inches(0.5), '行政書士に依頼', font_size=14, color=TEXT_MUTED, bold=True, alignment=PP_ALIGN.CENTER)
txb = add_text(s, Inches(1.2), Inches(4.2), Inches(6.5), Inches(1.2), '20万円超', font_size=44, color=TEXT_MUTED, bold=True, alignment=PP_ALIGN.CENTER)
# strikethrough
for p in txb.text_frame.paragraphs:
    for r in p.runs:
        r.font._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}strike'] = 'sngStrike'
add_text(s, Inches(1.2), Inches(5.6), Inches(6.5), Inches(1.0), '着手金5〜15万 + 成功報酬\nしかも毎回かかる', font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
# New price
shape_new = add_shape_rect(s, Inches(8.3), Inches(3.0), Inches(6.5), Inches(4.5), CORAL_BG, CORAL)
shape_new.adjustments[0] = 0.04
shape_new.line.width = Pt(2.5)
add_text(s, Inches(8.3), Inches(3.4), Inches(6.5), Inches(0.5), 'このコース', font_size=14, color=CORAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(8.3), Inches(4.2), Inches(6.5), Inches(1.2), '29,800円', font_size=48, color=CORAL, bold=True, alignment=PP_ALIGN.CENTER, font_name='Inter')
add_text(s, Inches(8.3), Inches(5.6), Inches(6.5), Inches(1.0), '買い切り・何度でも使える', font_size=16, color=TEXT_BODY, bold=True, alignment=PP_ALIGN.CENTER)

# Limited bonus
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), '今日だけの特典')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         '勉強会参加者限定特典', font_size=36, color=TEXT_DARK, bold=True)
add_highlight_box(s, Inches(1.2), Inches(3.2), Inches(13), Inches(1.5),
                  'パート6で実演した「補助金自動生成プログラム」\n（Claude Codeスキルパッケージ）', font_size=24)
add_bullet_list(s, Inches(1.2), Inches(5.5), Inches(13), Inches(2.0), [
    '通常はコースに含まれていない',
    '今日の勉強会参加者だけに追加でつけます',
], font_size=22)

# 3 roads
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, NAVY)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(0.5), Inches(2), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = CORAL
line.line.fill.background()
add_text(s, Inches(1.5), Inches(1.5), Inches(13), Inches(1.5), '3つの道', font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
# Card 1
c1 = add_shape_rect(s, Inches(1.5), Inches(3.5), Inches(3.8), Inches(3.5), RGBColor(0x28, 0x28, 0x42), RGBColor(0x3A, 0x3A, 0x55))
c1.adjustments[0] = 0.04
add_text(s, Inches(1.5), Inches(3.8), Inches(3.8), Inches(0.6), '1', font_size=28, color=TEXT_MUTED, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.8), Inches(3.8), Inches(1.5), '今日の話は\nなかったことにする', font_size=18, color=RGBColor(0x88, 0x88, 0x99), alignment=PP_ALIGN.CENTER)
# Card 2
c2 = add_shape_rect(s, Inches(6.1), Inches(3.5), Inches(3.8), Inches(3.5), RGBColor(0x28, 0x28, 0x42), RGBColor(0x3A, 0x3A, 0x55))
c2.adjustments[0] = 0.04
add_text(s, Inches(6.1), Inches(3.8), Inches(3.8), Inches(0.6), '2', font_size=28, color=TEXT_MUTED, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(6.1), Inches(4.8), Inches(3.8), Inches(1.5), '専門家に\n毎回依頼する', font_size=18, color=RGBColor(0x88, 0x88, 0x99), alignment=PP_ALIGN.CENTER)
# Card 3 (highlight)
c3 = add_shape_rect(s, Inches(10.7), Inches(3.5), Inches(3.8), Inches(3.5), RGBColor(0x3D, 0x1F, 0x1A), CORAL)
c3.adjustments[0] = 0.04
c3.line.width = Pt(2.5)
add_text(s, Inches(10.7), Inches(3.8), Inches(3.8), Inches(0.6), '3', font_size=28, color=CORAL_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(10.7), Inches(4.8), Inches(3.8), Inches(1.5), 'このコースで学んで\nAIと一緒にやる', font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# =====================================================================
# Q&A
# =====================================================================
section_slide('ご質問は\nありますか？', '', 'Q&A')

# Deadline
s = content_slide(WARM2)
add_label(s, Inches(1.2), Inches(0.6), 'DEADLINE')
add_text(s, Inches(1.2), Inches(1.2), Inches(13), Inches(1),
         'お申し込み締切', font_size=36, color=TEXT_DARK, bold=True)
add_highlight_box(s, Inches(2.0), Inches(3.5), Inches(12), Inches(1.8),
                  '2026年 4月13日（日）23:59', font_size=40)
add_text(s, Inches(1.2), Inches(6.5), Inches(13.5), Inches(1),
         '本日中にお申し込み方法をご案内いたします', font_size=22, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# Thank you
section_slide('ありがとうございました', 'Claude Code 導入勉強会')

# Save
output_path = '/Users/kazuhiroakutsu/claude-seminar-lp/slides/claude-code-seminar.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'Total slides: {len(prs.slides)}')
